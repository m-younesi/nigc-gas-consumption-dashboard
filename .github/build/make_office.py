# -*- coding: utf-8 -*-
"""
ساخت خروجی‌های اداری: پاورپوینت، اکسل و CSV.

تصاویر نمودارها از پیش با کروم از همان وب‌اپ رندر شده‌اند (پوشهٔ build/chart_png)
تا اسلایدها دقیقاً همان چیزی را نشان دهند که در وب می‌بینید.
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference, ScatterChart, Series
from openpyxl.formatting.rule import ColorScaleRule, DataBarRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
BUILD = ROOT / "build"
PNG = BUILD / "chart_png"
# نام فایل‌ها ASCII است: گیت‌هاب‌پیجز مسیرهای غیرلاتین را ۴۰۴ می‌دهد.
# نام فارسی هنگام دانلود، از صفت download در صفحه اعمال می‌شود.
OUT = ROOT / "webapp" / "downloads"
IMG = ROOT / "webapp" / "assets" / "img"

DATA = json.loads((BUILD / "dataset.json").read_text(encoding="utf-8"))
NAT = DATA["national"]

# متن و رنگ خروجی‌ها از پنل مدیریت می‌آید. اگر content.json نبود یا کلیدی
# نداشت، مقدار پیش‌فرضِ همین فایل به کار می‌رود تا ساخت هیچ‌وقت نشکند.
_CONTENT_PATH = ROOT / "webapp" / "content.json"
try:
    CONTENT = json.loads(_CONTENT_PATH.read_text(encoding="utf-8"))
except Exception:  # noqa: BLE001
    CONTENT = {}
OUTPUTS = CONTENT.get("outputs") or {}


def cfg(path, default=""):
    """مقدار را از outputs می‌خواند: cfg('office.pptx.title', 'پیش‌فرض')."""
    cur = OUTPUTS
    for part in path.split("."):
        if not isinstance(cur, dict) or part not in cur:
            return default
        cur = cur[part]
    return default if cur in (None, "") else cur


def hex_rgb(value, default):
    """«#rrggbb» را به RGBColor تبدیل می‌کند؛ ورودی نامعتبر را نادیده می‌گیرد."""
    text = str(value or "").strip().lstrip("#")
    if len(text) != 6:
        return default
    try:
        return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except ValueError:
        return default
PROV = DATA["provinces"]
H = DATA["headline"]
GK = ["low", "medium", "high", "very"]
GL = ["کم‌مصرف (الگو)", "مصرف متوسط", "پرمصرف", "بسیار پرمصرف"]
TIERS = [f"پلهٔ {i}" for i in range(1, 13)]

FA_DIGITS = str.maketrans("0123456789.", "۰۱۲۳۴۵۶۷۸۹٫")


def fa(value, decimals=1):
    """عدد را با ارقام فارسی برمی‌گرداند (برای متن اسلاید، نه سلول اکسل)."""
    if value is None:
        return "—"
    return f"{value:.{decimals}f}".translate(FA_DIGITS)


# ===================================================================== پاورپوینت
# پالت تیرهٔ پروژه (همان توکن‌های CSS)
PLANE = RGBColor(0x0A, 0x0E, 0x14)
SURF = RGBColor(0x14, 0x19, 0x22)
SURF2 = RGBColor(0x1B, 0x21, 0x30)
INK = RGBColor(0xF2, 0xF5, 0xF8)
INK2 = RGBColor(0xA8, 0xB3, 0xC2)
MUTED = RGBColor(0x6F, 0x7D, 0x8F)
SUBS = RGBColor(0x39, 0x87, 0xE5)
GAS = RGBColor(0xF0, 0x80, 0x40)
GROUPS = [RGBColor(0x8A, 0x4A, 0x12), RGBColor(0xC0, 0x6A, 0x1E),
          RGBColor(0xE0, 0x8A, 0x3A), RGBColor(0xF5, 0xB0, 0x6A)]

# فونت اسلاید و سلول‌های اکسل باید روی هر ویندوزی موجود باشد؛ استعداد فونت
# سیستمی نیست و پاورپوینت آن را با چیز نامعلومی جایگزین می‌کند.
# (متن داخل تصاویر نمودار همان استعداد است، چون در مرورگر رندر شده.)
FONT = "Tahoma"
W, HGT = Inches(13.333), Inches(7.5)


def set_rtl(frame):
    """جهت راست‌به‌چپ را روی همهٔ پاراگراف‌های یک کادر متن اعمال می‌کند."""
    for para in frame.paragraphs:
        para.alignment = PP_ALIGN.RIGHT
        pPr = para._p.get_or_add_pPr()
        pPr.set("rtl", "1")
        pPr.set("algn", "r")


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=INK,
            align=PP_ALIGN.RIGHT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        para.line_spacing = spacing
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = FONT
    set_rtl(tf)
    for para in tf.paragraphs:
        para.alignment = align
    return box


def bg(slide, color=PLANE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.06
        except Exception:
            pass
    return shape


def slide_header(slide, kicker, title, subtitle=None):
    textbox(slide, Inches(0.62), Inches(0.38), Inches(12.1), Inches(0.34),
            kicker, size=12, bold=True, color=GAS)
    textbox(slide, Inches(0.62), Inches(0.72), Inches(12.1), Inches(0.72),
            title, size=30, bold=True, color=INK)
    if subtitle:
        textbox(slide, Inches(0.62), Inches(1.45), Inches(12.1), Inches(0.55),
                subtitle, size=13, color=INK2)


def add_picture_fit(slide, path, x, y, w, h):
    """تصویر را با حفظ نسبت، درون کادر داده‌شده وسط‌چین می‌کند."""
    from PIL import Image as PILImage
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(str(path), int(x + (w - nw) / 2),
                                    int(y + (h - nh) / 2), nw, nh)


def build_pptx():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, HGT
    blank = prs.slide_layouts[6]

    # ------------------------------------------------- ۱ جلد
    s = prs.slides.add_slide(blank)
    bg(s)
    hero = IMG / "onshore.jpg"
    if hero.exists():
        pic = s.shapes.add_picture(str(hero), 0, 0, W, HGT)
        pic.crop_top = pic.crop_bottom = 0.06
    veil = rect(s, 0, 0, W, HGT, PLANE, radius=False)
    veil.fill.transparency = 0.30
    band = rect(s, 0, Inches(3.5), W, Inches(4.0), PLANE, radius=False)
    band.fill.transparency = 0.12

    logo = IMG / "nigc-logo.svg"
    org_line = "%s — %s" % (cfg("shared.org", "شرکت ملی گاز ایران"),
                            cfg("shared.department",
                                "مدیریت گازرسانی، امور تعرفه‌ها و قراردادها"))
    accent = hex_rgb(cfg("office.pptx.accent"), GAS)
    textbox(s, Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.4),
            org_line, size=13, color=INK2)
    textbox(s, Inches(0.8), Inches(4.15), Inches(11.7), Inches(1.5),
            cfg("office.pptx.title", "آناتومی یک زمستان"),
            size=48, bold=True, color=INK)
    textbox(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.8),
            cfg("office.pptx.subtitle", "چه کسی گاز ایران را می‌سوزاند؟"),
            size=32, bold=True, color=accent)
    textbox(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9),
            cfg("office.pptx.cover_note",
                "تحلیل مصرف گاز بخش خانگی سال ۱۴۰۴ به تفکیک پلهٔ تعرفه و استان")
            + "\n"
            + "%s · گزارش %s" % (cfg("shared.scope",
                                     "۳۱ استان · ۱۲ پلهٔ دورهٔ سرد · ۴ پلهٔ دورهٔ گرم"),
                                 cfg("shared.report_date", "۱۴۰۵/۰۵/۱۸")),
            size=13, color=INK2)

    # ------------------------------------------------- ۲ خلاصه در چهار عدد
    s = prs.slides.add_slide(blank)
    bg(s)
    slide_header(s, "خلاصهٔ مدیریتی", "چهار عددی که کل ماجرا را می‌گوید")
    stats = [
        (fa(H["low_count"]) + "٪", "مشترکین گروه کم‌مصرف",
         f"که {fa(H['low_cons'])}٪ از گاز را مصرف می‌کنند", SUBS),
        (fa(H["over_pattern"]) + "٪", "گاز سوخته بالاتر از الگو",
         f"توسط تنها {fa(H['over_pattern_count'])}٪ از مشترکین", GAS),
        (fa(H["top10"]) + "٪", "سهم ۱۰٪ پرمصرف‌ترین",
         f"نیمهٔ کم‌مصرف تنها {fa(H['bottom50'])}٪ را می‌برد", GROUPS[2]),
        (fa(H["tier12_multiplier"]) + "×", "شدت مصرف مشترک پلهٔ ۱۲",
         "معادل حدود ۱۱ مشترک پلهٔ ۱", GROUPS[3]),
    ]
    cw, gap = Inches(2.92), Inches(0.22)
    for i, (val, key, desc, color) in enumerate(stats):
        x = W - Inches(0.62) - cw - i * (cw + gap)
        rect(s, x, Inches(2.35), cw, Inches(2.5), SURF)
        rect(s, x + cw - Inches(0.06), Inches(2.35), Inches(0.06), Inches(2.5), color, radius=False)
        textbox(s, x + Inches(0.2), Inches(2.6), cw - Inches(0.45), Inches(0.9),
                val, size=40, bold=True, color=INK)
        textbox(s, x + Inches(0.2), Inches(3.55), cw - Inches(0.45), Inches(0.5),
                key, size=13, bold=True, color=INK)
        textbox(s, x + Inches(0.2), Inches(4.05), cw - Inches(0.45), Inches(0.7),
                desc, size=11, color=MUTED)

    textbox(s, Inches(0.62), Inches(5.2), Inches(12.1), Inches(1.6),
            "دو توزیع موازی وجود دارد: یکی برای تعداد مشترکین و یکی برای حجم گاز. "
            "اگر بر هم منطبق بودند، مصرف کاملاً برابر بود. فاصلهٔ میان آن‌ها همان چیزی "
            "است که این گزارش می‌سنجد.", size=15, color=INK2)

    # ------------------------------------------------- اسلایدهای نموداری
    chart_slides = [
        ("butterfly.png", "نردبان پله‌ها",
         "مشترک در برابر مصرف، پله به پله",
         f"پلهٔ ۱ با {fa(NAT['h1_count'][0])}٪ از مشترکین، {fa(NAT['h1_cons'][0])}٪ از گاز را "
         f"می‌برد؛ پلهٔ ۱۲ با تنها {fa(NAT['h1_count'][11])}٪ از مشترکین، "
         f"{fa(NAT['h1_cons'][11])}٪ را."),
        ("intensity.png", "ضریب شدت",
         "هر پله چند برابر مشترک متوسط گاز می‌سوزاند",
         f"شدت مصرف از {fa(NAT['h1_intensity'][0])}× در پلهٔ ۱ تا "
         f"{fa(NAT['h1_intensity'][11])}× در پلهٔ ۱۲ بالا می‌رود."),
        ("sankey.png", "جریان مصرف",
         "از صد مشترک تا صد واحد گاز",
         f"گروه بسیار پرمصرف با {fa(NAT['h1_count_g']['very'])}٪ مشترک، "
         f"{fa(NAT['h1_cons_g']['very'])}٪ گاز را می‌برد."),
        ("lorenz.png", "نابرابری مصرف",
         "منحنی لورنز و ضریب جینی",
         f"ضریب جینی کشوری {fa(H['gini'], 3)} است — کران پایین نابرابری واقعی، چون "
         f"نابرابری درون هر پله در داده دیده نمی‌شود."),
        ("map.png", "جغرافیای مصرف",
         "سهم مصرف مازاد بر الگو، به تفکیک استان",
         f"بیشترین در {DATA['rankings']['over_pattern'][0]['province']} "
         f"({fa(DATA['rankings']['over_pattern'][0]['value'])}٪) و کمترین در "
         f"{DATA['rankings']['over_pattern'][-1]['province']} "
         f"({fa(DATA['rankings']['over_pattern'][-1]['value'])}٪)."),
        ("heat.png", "ماتریس استان × پله",
         "الگوی توزیع مصرف در هر استان",
         "ردیف‌هایی که روشنی‌شان به چپ کشیده شده، دنبالهٔ پرمصرف سنگینی دارند."),
        ("climate.png", "اقلیم در برابر رفتار",
         "دما چقدر از تفاوت استان‌ها را توضیح می‌دهد؟",
         f"دما تنها {fa(DATA['climate']['over_pattern']['r2'] * 100, 0)}٪ از پراکندگی را "
         f"توضیح می‌دهد؛ بقیه‌اش رفتار، عایق‌بندی و رطوبت است."),
        ("season.png", "دو فصل",
         "فاصلهٔ پرمصرف‌ترین گروه تا مشترک متوسط",
         f"از {fa(NAT['h2_intensity_g']['very'])}× در دورهٔ گرم به "
         f"{fa(NAT['h1_intensity_g']['very'])}× در دورهٔ سرد."),
    ]

    # اگر در پنل خاموش شده باشد، اسلایدهای نمودار ساخته نمی‌شوند
    if not cfg("office.pptx.include_charts", True):
        chart_slides = []

    for fname, kicker, title, note in chart_slides:
        path = PNG / fname
        if not path.exists():
            print("  ! تصویر نمودار نیست:", fname)
            continue
        s = prs.slides.add_slide(blank)
        bg(s)
        slide_header(s, kicker, title)
        card = rect(s, Inches(0.62), Inches(1.55), Inches(12.1), Inches(4.6), SURF)
        add_picture_fit(s, path, Inches(0.75), Inches(1.68),
                        Inches(11.84), Inches(4.34))
        rect(s, Inches(0.62), Inches(6.3), Inches(12.1), Inches(0.78), SURF2)
        rect(s, Inches(12.66), Inches(6.3), Inches(0.06), Inches(0.78), GAS, radius=False)
        textbox(s, Inches(0.85), Inches(6.44), Inches(11.6), Inches(0.6),
                note, size=13, color=INK2)

    # ------------------------------------------------- رتبه‌بندی استان‌ها
    s = prs.slides.add_slide(blank)
    bg(s)
    slide_header(s, "رتبه‌بندی", "پرمصرف‌ترین و کم‌مصرف‌ترین استان‌ها نسبت به الگو")
    ranked = DATA["rankings"]["over_pattern"]
    top, bottom = ranked[:8], ranked[-8:]
    for col, (rows, label, color) in enumerate(
            [(top, "بیشترین مصرف مازاد بر الگو", GAS),
             (list(reversed(bottom)), "کمترین مصرف مازاد بر الگو", SUBS)]):
        x = W - Inches(0.62) - Inches(6.0) - col * Inches(6.3)
        textbox(s, x, Inches(1.95), Inches(6.0), Inches(0.4), label,
                size=15, bold=True, color=INK)
        mx = max(r["value"] for r in ranked)
        for i, r in enumerate(rows):
            y = Inches(2.5) + i * Inches(0.56)
            rect(s, x, y, Inches(6.0), Inches(0.46), SURF)
            textbox(s, x + Inches(3.0), y + Inches(0.05), Inches(2.85), Inches(0.36),
                    r["province"], size=12, bold=True, color=INK)
            bw = Inches(2.4) * (r["value"] / mx)
            rect(s, x + Inches(2.8) - bw, y + Inches(0.16), bw, Inches(0.14),
                 color, radius=False)
            textbox(s, x + Inches(0.15), y + Inches(0.05), Inches(0.9), Inches(0.36),
                    fa(r["value"]) + "٪", size=12, color=INK2, align=PP_ALIGN.LEFT)

    # ------------------------------------------------- روش‌شناسی
    s = prs.slides.add_slide(blank)
    bg(s)
    slide_header(s, "پیوست", "روش‌شناسی و محدودیت‌ها")
    left_txt = (
        "منبع داده\n"
        "چهار جدول منتشرشدهٔ مدیریت گازرسانی — امور تعرفه‌ها و قراردادها: "
        "درصد مصرف و درصد تعداد مشترکین بخش خانگی، برای نیمهٔ اول (دورهٔ سرد، ۱۲ پله) "
        "و نیمهٔ دوم (دورهٔ گرم، ۴ پله) سال ۱۴۰۴.\n\n"
        "اعتبارسنجی\n"
        "جمع هر ردیف باید به ۱۰۰ برسد و مجموع پله‌های هر گروه باید با ردیف‌های تجمیعی "
        "جدول مبدأ بخواند. در همین بررسی یک خطای انتقالی در ردیف آذربایجان شرقی "
        "شناسایی و اصلاح شد."
    )
    right_txt = (
        "محدودیت‌ها\n"
        "• ارقام مبدأ به یک رقم اعشار گرد شده‌اند؛ هر ردیف برای محاسبه به مجموع ۱۰۰ "
        "نرمال شده است.\n"
        "• جینی محاسبه‌شده کران پایین نابرابری واقعی است؛ نابرابری درون هر پله دیده نمی‌شود.\n"
        "• داده «درصد» است، نه حجم مطلق. سهم هر استان از کل گاز کشور از این جدول‌ها "
        "قابل استخراج نیست.\n"
        "• پله‌های دو نیم‌سال هم‌ارز نیستند؛ هیچ سهم خامی بین دو دوره مقایسه نشده است."
    )
    rect(s, Inches(6.95), Inches(2.05), Inches(5.77), Inches(4.4), SURF)
    textbox(s, Inches(7.2), Inches(2.3), Inches(5.3), Inches(4.0), left_txt,
            size=12, color=INK2)
    rect(s, Inches(0.62), Inches(2.05), Inches(6.05), Inches(4.4), SURF)
    textbox(s, Inches(0.87), Inches(2.3), Inches(5.6), Inches(4.0), right_txt,
            size=12, color=INK2)

    OUT.mkdir(parents=True, exist_ok=True)
    dest = OUT / "gas-presentation-1404.pptx"
    prs.save(dest)
    print(f"PPTX  {dest.name}  {len(prs.slides._sldIdLst)} اسلاید  "
          f"{dest.stat().st_size/1024:.0f} KB")


# ========================================================================= اکسل
THIN = Side(style="thin", color="D8D7D0")
HEAD_FILL = PatternFill("solid", fgColor="1B2130")
TOT_FILL = PatternFill("solid", fgColor="EFEDE6")


def style_header(ws, row=1, cols=None):
    for cell in ws[row]:
        if cols and cell.column > cols:
            break
        cell.font = Font(name=FONT, bold=True, size=10, color="FFFFFF")
        cell.fill = HEAD_FILL
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = Border(bottom=THIN)
    ws.row_dimensions[row].height = 34


def autosize(ws, min_w=9, max_w=26):
    for col in ws.columns:
        letter = get_column_letter(col[0].column)
        longest = max((len(str(c.value)) for c in col if c.value is not None), default=0)
        ws.column_dimensions[letter].width = max(min_w, min(max_w, longest + 3))


def build_xlsx():
    wb = Workbook()

    # ---------------------------------------------------- برگهٔ راهنما
    ws = wb.active
    ws.title = "راهنما"
    ws.sheet_view.rightToLeft = True
    guide = [
        (cfg("office.xlsx.title",
             "تحلیل مصرف گاز خانگی به تفکیک پله — سال ۱۴۰۴"), ""),
        ("", ""),
        ("منبع", "%s — %s" % (cfg("shared.department",
                                  "مدیریت گازرسانی، امور تعرفه‌ها و قراردادها"),
                              cfg("shared.org", "شرکت ملی گاز ایران"))),
        ("تاریخ گزارش مبدأ", cfg("shared.report_date", "۱۴۰۵/۰۵/۱۸")),
        ("دامنه", cfg("shared.scope",
                      "۳۱ استان × ۱۲ پلهٔ دورهٔ سرد × ۴ پلهٔ دورهٔ گرم")),
        ("", ""),
        ("برگه‌ها", ""),
        ("شاخص‌ها", "شاخص‌های محاسبه‌شده برای هر استان + نمودار میله‌ای"),
        ("مصرف - دورهٔ سرد", "درصد مصرف در هر یک از ۱۲ پله"),
        ("مشترکین - دورهٔ سرد", "درصد تعداد مشترکین در هر یک از ۱۲ پله"),
        ("دورهٔ گرم", "درصد مصرف و تعداد مشترکین در ۴ پله"),
        ("گروه‌ها", "تجمیع چهار گروه مصرف در هر دو دوره"),
        ("منحنی لورنز", "نقاط منحنی لورنز کشوری + نمودار"),
        ("", ""),
        ("تعریف شاخص‌ها", ""),
        ("ضریب شدت", "سهم مصرف پله ÷ سهم مشترکین پله"),
        ("مازاد بر الگو", "۱۰۰ منهای سهم مصرف گروه کم‌مصرف (پله‌های ۱ تا ۳)"),
        ("ضریب جینی", "۱ منهای دو برابر مساحت زیر منحنی لورنز — کران پایین نابرابری"),
        ("انحراف اقلیمی", "اختلاف مصرف مازاد بر الگو با پیش‌بینی رگرسیون بر حسب دما"),
        ("", ""),
        ("هشدار", "ارقام مبدأ گرد شده‌اند و هر ردیف به مجموع ۱۰۰ نرمال شده است."),
        ("هشدار", "پله‌های دو نیم‌سال هم‌ارز نیستند؛ سهم‌های خام دو دوره مقایسه‌شدنی نیستند."),
    ]
    for r, (a, b) in enumerate(guide, start=1):
        ws.cell(r, 1, a).font = Font(name=FONT, bold=(b == "" and a != ""), size=11)
        ws.cell(r, 2, b).font = Font(name=FONT, size=10)
        ws.cell(r, 2).alignment = Alignment(wrap_text=True, vertical="top")
    ws.cell(1, 1).font = Font(name=FONT, bold=True, size=15, color="B04C0F")
    ws.column_dimensions["A"].width = 26
    ws.column_dimensions["B"].width = 82
    ws.freeze_panes = "A2"

    # ---------------------------------------------------- شاخص‌ها
    ws = wb.create_sheet("شاخص‌ها")
    ws.sheet_view.rightToLeft = True
    head = ["استان", "مازاد بر الگو ٪", "سهم پرمصرف‌ها ٪", "سهم ۱۰٪ بالا ٪",
            "سهم ۵۰٪ پایین ٪", "جینی (سرد)", "جینی (گرم)", "پلهٔ میانه",
            "مشترکین الگو ٪", "مصرف الگو ٪", "مصرف بسیار پرمصرف ٪",
            "شدت بسیار پرمصرف", "میانگین دما °C", "انتظار اقلیمی ٪", "انحراف اقلیمی"]
    ws.append(head)
    rows = sorted(PROV, key=lambda r: -r["h1_over_pattern"])
    for r in rows:
        ws.append([
            r["province"], r["h1_over_pattern"], r["h1_tail_cons"], r["h1_top10"],
            r["h1_bottom50"], r["h1_gini"], r["h2_gini"], r["h1_median_tier"],
            r["h1_count_g"]["low"], r["h1_cons_g"]["low"], r["h1_cons_g"]["very"],
            r["h1_intensity_g"]["very"], r["temp"], r["climate_expected"],
            r["climate_residual"],
        ])
    ws.append(["کل کشور", NAT["h1_over_pattern"], NAT["h1_tail_cons"], NAT["h1_top10"],
               NAT["h1_bottom50"], NAT["h1_gini"], NAT["h2_gini"], NAT["h1_median_tier"],
               NAT["h1_count_g"]["low"], NAT["h1_cons_g"]["low"], NAT["h1_cons_g"]["very"],
               NAT["h1_intensity_g"]["very"], None, None, None])
    style_header(ws)
    last = ws.max_row
    for row in ws.iter_rows(min_row=2, max_row=last):
        for c in row:
            c.font = Font(name=FONT, size=10, bold=(row[0].row == last))
            c.alignment = Alignment(horizontal="center")
            c.border = Border(bottom=THIN)
            if isinstance(c.value, float):
                c.number_format = "0.000" if c.column in (6, 7) else "0.0"
        row[0].alignment = Alignment(horizontal="right")
        if row[0].row == last:
            for c in row:
                c.fill = TOT_FILL
    ws.freeze_panes = "B2"
    autosize(ws)

    ws.conditional_formatting.add(f"B2:B{last-1}",
        ColorScaleRule(start_type="min", start_color="CDE2FB",
                       end_type="max", end_color="0D366B"))
    ws.conditional_formatting.add(f"O2:O{last-1}",
        ColorScaleRule(start_type="min", start_color="2A78D6",
                       mid_type="num", mid_value=0, mid_color="F0EFEC",
                       end_type="max", end_color="E34948"))
    ws.conditional_formatting.add(f"F2:F{last-1}",
        DataBarRule(start_type="min", end_type="max", color="D97324"))

    chart = BarChart()
    chart.type = "bar"
    chart.title = "سهم مصرف مازاد بر الگو به تفکیک استان"
    chart.y_axis.title = "درصد"
    chart.height, chart.width = 18, 24
    chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=last - 1), titles_from_data=True)
    chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=last - 1))
    chart.legend = None
    ws.add_chart(chart, "Q2")

    # ---------------------------------------------------- جدول‌های خام ۱۲ پله
    for key, title in [("h1_cons", "مصرف - دورهٔ سرد"), ("h1_count", "مشترکین - دورهٔ سرد")]:
        ws = wb.create_sheet(title)
        ws.sheet_view.rightToLeft = True
        ws.append(["استان"] + TIERS + ["جمع"])
        for r in PROV:
            ws.append([r["province"]] + r[key] + [round(sum(r[key]), 1)])
        ws.append(["کل کشور"] + NAT[key] + [round(sum(NAT[key]), 1)])
        style_header(ws)
        last = ws.max_row
        for row in ws.iter_rows(min_row=2, max_row=last):
            for c in row:
                c.font = Font(name=FONT, size=10, bold=(row[0].row == last))
                c.alignment = Alignment(horizontal="center")
                c.border = Border(bottom=THIN)
                if isinstance(c.value, float):
                    c.number_format = "0.0"
            row[0].alignment = Alignment(horizontal="right")
            if row[0].row == last:
                for c in row:
                    c.fill = TOT_FILL
        ws.freeze_panes = "B2"
        autosize(ws, min_w=7, max_w=22)
        ws.conditional_formatting.add(f"B2:M{last}",
            ColorScaleRule(start_type="num", start_value=0, start_color="F5F8FD",
                           mid_type="num", mid_value=8, mid_color="6DA7EC",
                           end_type="num", end_value=35, end_color="0D366B"))

    # ---------------------------------------------------- دورهٔ گرم
    ws = wb.create_sheet("دورهٔ گرم")
    ws.sheet_view.rightToLeft = True
    ws.append(["استان", "مصرف پلهٔ ۱ ٪", "مصرف پلهٔ ۲ ٪", "مصرف پلهٔ ۳ ٪", "مصرف پلهٔ ۴ ٪",
               "مشترکین پلهٔ ۱ ٪", "مشترکین پلهٔ ۲ ٪", "مشترکین پلهٔ ۳ ٪", "مشترکین پلهٔ ۴ ٪",
               "جینی"])
    for r in PROV + [NAT]:
        ws.append([r["province"]] + r["h2_cons"] + r["h2_count"] + [r["h2_gini"]])
    style_header(ws)
    last = ws.max_row
    for row in ws.iter_rows(min_row=2, max_row=last):
        for c in row:
            c.font = Font(name=FONT, size=10, bold=(row[0].row == last))
            c.alignment = Alignment(horizontal="center")
            c.border = Border(bottom=THIN)
            if isinstance(c.value, float):
                c.number_format = "0.000" if c.column == 10 else "0.0"
        row[0].alignment = Alignment(horizontal="right")
    ws.freeze_panes = "B2"
    autosize(ws, min_w=11)

    # ---------------------------------------------------- گروه‌ها
    ws = wb.create_sheet("گروه‌ها")
    ws.sheet_view.rightToLeft = True
    ws.append(["استان"] +
              [f"مصرف {g} ٪ (سرد)" for g in GL] +
              [f"مشترکین {g} ٪ (سرد)" for g in GL] +
              [f"شدت {g} (سرد)" for g in GL])
    for r in PROV + [NAT]:
        ws.append([r["province"]] +
                  [r["h1_cons_g"][k] for k in GK] +
                  [r["h1_count_g"][k] for k in GK] +
                  [r["h1_intensity_g"][k] for k in GK])
    style_header(ws)
    last = ws.max_row
    for row in ws.iter_rows(min_row=2, max_row=last):
        for c in row:
            c.font = Font(name=FONT, size=10, bold=(row[0].row == last))
            c.alignment = Alignment(horizontal="center")
            c.border = Border(bottom=THIN)
            if isinstance(c.value, float):
                c.number_format = "0.00" if c.column > 9 else "0.0"
        row[0].alignment = Alignment(horizontal="right")
    ws.freeze_panes = "B2"
    autosize(ws, min_w=13)

    # ---------------------------------------------------- لورنز
    ws = wb.create_sheet("منحنی لورنز")
    ws.sheet_view.rightToLeft = True
    ws.append(["سهم تجمعی مشترکین ٪", "سهم تجمعی مصرف ٪", "خط برابری کامل"])
    for px, qy in NAT["h1_lorenz"]:
        ws.append([px, qy, px])
    style_header(ws)
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
        for c in row:
            c.font = Font(name=FONT, size=10)
            c.alignment = Alignment(horizontal="center")
            c.number_format = "0.0"
    autosize(ws, min_w=20)

    lc = ScatterChart()
    lc.title = "منحنی لورنز مصرف گاز خانگی — کل کشور، دورهٔ سرد ۱۴۰۴"
    lc.x_axis.title = "درصد تجمعی مشترکین"
    lc.y_axis.title = "درصد تجمعی مصرف"
    lc.height, lc.width = 12, 16
    xref = Reference(ws, min_col=1, min_row=2, max_row=ws.max_row)
    for col, name in [(2, "مصرف واقعی"), (3, "برابری کامل")]:
        ser = Series(Reference(ws, min_col=col, min_row=1, max_row=ws.max_row),
                     xref, title_from_data=True)
        ser.marker.symbol = "none"
        lc.series.append(ser)
    ws.add_chart(lc, "E2")

    OUT.mkdir(parents=True, exist_ok=True)
    dest = OUT / "gas-data-analysis-1404.xlsx"
    wb.save(dest)
    print(f"XLSX  {dest.name}  {len(wb.sheetnames)} برگه  {dest.stat().st_size/1024:.0f} KB")


# =========================================================================== CSV
def build_csv():
    OUT.mkdir(parents=True, exist_ok=True)
    dest = OUT / "gas-raw-data-1404.csv"
    head = (["province"] +
            [f"h1_cons_p{i}" for i in range(1, 13)] +
            [f"h1_count_p{i}" for i in range(1, 13)] +
            [f"h2_cons_p{i}" for i in range(1, 5)] +
            [f"h2_count_p{i}" for i in range(1, 5)] +
            ["h1_over_pattern", "h1_tail_cons", "h1_top10", "h1_bottom50",
             "h1_gini", "h2_gini", "h1_median_tier", "temp",
             "climate_expected", "climate_residual"])
    with dest.open("w", encoding="utf-8-sig", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(head)
        for r in PROV + [NAT]:
            w.writerow([r["province"]] + r["h1_cons"] + r["h1_count"] +
                       r["h2_cons"] + r["h2_count"] +
                       [r["h1_over_pattern"], r["h1_tail_cons"], r["h1_top10"],
                        r["h1_bottom50"], r["h1_gini"], r["h2_gini"],
                        r["h1_median_tier"], r["temp"],
                        r["climate_expected"], r["climate_residual"]])
    print(f"CSV   {dest.name}  {dest.stat().st_size/1024:.0f} KB")


if __name__ == "__main__":
    build_csv()
    build_xlsx()
    build_pptx()
